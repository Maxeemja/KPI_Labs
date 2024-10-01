from pptx import Presentation
from pptx.util import Inches

# Створення презентації
prs = Presentation()

# Слайд 1: Титульний слайд
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Аналіз логістичної системи ТОВ \"КБ Альфа\""
subtitle.text = ""

# Слайд 2: Загальна характеристика
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "Загальна характеристика"
content.text = "- Виробник та постачальник агротехніки і обладнання\n" \
               "- Функціонує з 2011 року, реорганізація у 2015-2017 рр.\n" \
               "- Стабільні фінансові показники з невисокою рентабельністю до 2022 р."

# Слайд 3: Маркетингова стратегія
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]
title.text = "Маркетингова стратегія"
content.text = "- Зосереджена на комунікаційній політиці\n" \
               "- Участь у виставках, реклама, стимулювання збуту\n" \
               "- Орієнтація на ринок В2В"

# Слайд 4: Логістична система
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]
title.text = "Логістична система"
content.text = "- Закупівлі у ключових постачальників\n" \
               "- Власне виробництво вузького, але гармонійного асортименту\n" \
               "- Канали розподілу через прямі продажі, оптовиків, роздрібну торгівлю"

# Слайд 5: Проблеми та шляхи вирішення
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]
title.text = "Проблеми та шляхи вирішення"
content.text = "- Евакуація з Бердянська до Дніпра через війну\n" \
               "- Втрата частини активів та ринків збуту на окупованих територіях\n" \
               "- Рекомендації:\n" \
               "    - Відкрити представництво у м. Тернопіль для охоплення Західного регіону\n" \
               "    - Активне використання прямих продажів через комівояжерів\n" \
               "    - Посилення присутності у перспективних неохоплених регіонах\n" \
               "    - Мінімізація ризиків завдяки віддаленості та близькості до ЄС"

# Збереження презентації
prs.save('аналіз_логістичної_системи_ТОВ_КБ_Альфа.pptx')
